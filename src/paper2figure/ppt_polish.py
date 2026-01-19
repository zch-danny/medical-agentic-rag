"""
PPTPolish - PPT 美化和优化模块

对已有 PPT 进行美化、内容优化、布局调整
"""

import json
import os
import re
from dataclasses import dataclass
from enum import Enum
from pathlib import Path
from typing import Dict, List, Optional, Tuple, Union

from loguru import logger
from openai import OpenAI
from tenacity import retry, stop_after_attempt, wait_exponential


class PolishMode(str, Enum):
    """美化模式"""
    CONTENT = "content"      # 仅优化内容文字
    STYLE = "style"          # 仅调整样式
    FULL = "full"            # 完整优化（内容+样式）
    LAYOUT = "layout"        # 优化布局


@dataclass
class PolishResult:
    """美化结果"""
    original_path: str
    output_path: str
    changes: List[str]       # 修改记录
    suggestions: List[str]   # 优化建议


class PPTPolish:
    """
    PPT 美化器

    功能：
    1. 内容优化：精炼文字、统一风格
    2. 样式美化：调整颜色、字体、间距
    3. 布局优化：重新排版、对齐元素
    4. 一致性检查：统一格式和样式
    """

    CONTENT_POLISH_PROMPT = """你是一个专业的 PPT 内容优化师。请优化以下幻灯片内容：

当前内容：
{content}

优化要求：
1. 标题简洁有力（不超过 15 字）
2. 要点精炼，使用动词开头
3. 每条要点不超过 20 字
4. 保持专业术语的准确性
5. 添加适当的过渡语

请以 JSON 格式输出优化后的内容：
```json
{{
    "slides": [
        {{
            "original_title": "原标题",
            "new_title": "优化后标题",
            "original_points": ["原要点1", "原要点2"],
            "new_points": ["优化后要点1", "优化后要点2"],
            "suggestions": ["建议1"]
        }}
    ],
    "global_suggestions": ["全局建议"]
}}
```

只输出 JSON。"""

    # 预设配色方案
    COLOR_SCHEMES = {
        "academic_blue": {
            "name": "学术蓝",
            "primary": (26, 58, 92),
            "secondary": (74, 144, 217),
            "accent": (230, 126, 34),
            "text": (44, 62, 80),
            "background": (255, 255, 255),
        },
        "modern_green": {
            "name": "现代绿",
            "primary": (0, 128, 128),
            "secondary": (0, 184, 148),
            "accent": (255, 107, 107),
            "text": (52, 58, 64),
            "background": (248, 249, 250),
        },
        "elegant_purple": {
            "name": "优雅紫",
            "primary": (102, 92, 255),
            "secondary": (0, 201, 167),
            "accent": (255, 102, 144),
            "text": (51, 51, 51),
            "background": (255, 255, 255),
        },
        "business_navy": {
            "name": "商务蓝",
            "primary": (13, 71, 161),
            "secondary": (66, 165, 245),
            "accent": (255, 179, 0),
            "text": (33, 33, 33),
            "background": (255, 255, 255),
        },
        "warm_orange": {
            "name": "温暖橙",
            "primary": (230, 126, 34),
            "secondary": (241, 196, 15),
            "accent": (52, 73, 94),
            "text": (44, 62, 80),
            "background": (255, 255, 255),
        },
        "minimal_gray": {
            "name": "极简灰",
            "primary": (52, 73, 94),
            "secondary": (149, 165, 166),
            "accent": (231, 76, 60),
            "text": (44, 62, 80),
            "background": (236, 240, 241),
        },
    }

    # 字体方案
    FONT_SCHEMES = {
        "professional": {
            "name": "专业",
            "title_font": "Microsoft YaHei UI",
            "body_font": "Microsoft YaHei",
            "title_size": 32,
            "body_size": 20,
        },
        "elegant": {
            "name": "优雅",
            "title_font": "SimHei",
            "body_font": "SimSun",
            "title_size": 36,
            "body_size": 18,
        },
        "modern": {
            "name": "现代",
            "title_font": "Arial",
            "body_font": "Calibri",
            "title_size": 34,
            "body_size": 22,
        },
    }

    def __init__(
        self,
        api_key: Optional[str] = None,
        base_url: Optional[str] = None,
        model: Optional[str] = None,
    ):
        self.api_key = api_key or os.getenv("LLM_API_KEY")
        self.base_url = base_url or os.getenv("LLM_BASE_URL", "https://api.deepseek.com")
        self.model = model or os.getenv("LLM_MODEL", "deepseek-chat")

        if self.api_key:
            self.client = OpenAI(api_key=self.api_key, base_url=self.base_url)
        else:
            self.client = None

        logger.info("PPTPolish 初始化完成")

    def _extract_json(self, response: str) -> dict:
        """从 LLM 响应中提取 JSON"""
        try:
            return json.loads(response)
        except json.JSONDecodeError:
            pass

        pattern = r'```json\s*([\s\S]*?)```'
        match = re.search(pattern, response)
        if match:
            try:
                return json.loads(match.group(1))
            except json.JSONDecodeError:
                pass

        pattern = r'\{[\s\S]*\}'
        match = re.search(pattern, response)
        if match:
            try:
                return json.loads(match.group(0))
            except json.JSONDecodeError:
                pass

        raise ValueError("无法提取 JSON")

    @retry(
        stop=stop_after_attempt(3),
        wait=wait_exponential(multiplier=1, min=2, max=10),
        reraise=True,
    )
    def _call_llm(self, prompt: str) -> str:
        """调用 LLM"""
        if not self.client:
            raise ValueError("LLM 客户端未初始化")

        response = self.client.chat.completions.create(
            model=self.model,
            messages=[
                {"role": "system", "content": "你是一个专业的 PPT 内容优化师。"},
                {"role": "user", "content": prompt},
            ],
            temperature=0.3,
            max_tokens=4096,
        )
        return response.choices[0].message.content

    def _load_pptx(self, pptx_path: Union[str, Path]):
        """加载 PPTX 文件"""
        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError("python-pptx 未安装")

        return Presentation(str(pptx_path))

    def _extract_text_content(self, prs) -> List[Dict]:
        """提取 PPT 文字内容"""
        slides_content = []

        for i, slide in enumerate(prs.slides):
            slide_data = {
                "index": i,
                "title": "",
                "points": [],
            }

            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue

                text = shape.text_frame.text.strip()
                if not text:
                    continue

                # 简单判断是否为标题（通常字体较大或位置靠上）
                if shape.top and shape.top.inches < 2:
                    if not slide_data["title"]:
                        slide_data["title"] = text
                else:
                    # 分割多行文本
                    for para in shape.text_frame.paragraphs:
                        para_text = para.text.strip()
                        if para_text and para_text != slide_data["title"]:
                            slide_data["points"].append(para_text)

            slides_content.append(slide_data)

        return slides_content

    def polish_content(
        self,
        pptx_path: Union[str, Path],
        output_path: Optional[Union[str, Path]] = None,
    ) -> PolishResult:
        """
        优化 PPT 内容

        使用 LLM 优化文字内容，使其更加专业精炼

        Args:
            pptx_path: 输入 PPT 路径
            output_path: 输出路径（默认覆盖原文件）

        Returns:
            PolishResult: 优化结果
        """
        pptx_path = Path(pptx_path)
        if not pptx_path.exists():
            raise FileNotFoundError(f"文件不存在: {pptx_path}")

        if not output_path:
            output_path = pptx_path.parent / f"{pptx_path.stem}_polished.pptx"

        # 加载 PPT
        prs = self._load_pptx(pptx_path)

        # 提取内容
        content = self._extract_text_content(prs)

        # 调用 LLM 优化
        if self.client:
            prompt = self.CONTENT_POLISH_PROMPT.format(content=json.dumps(content, ensure_ascii=False, indent=2))
            response = self._call_llm(prompt)
            optimized = self._extract_json(response)
        else:
            optimized = {"slides": [], "global_suggestions": ["未配置 LLM，跳过内容优化"]}

        changes = []
        suggestions = optimized.get("global_suggestions", [])

        # 应用优化（简化版：记录建议但不直接修改）
        for opt_slide in optimized.get("slides", []):
            if opt_slide.get("new_title") != opt_slide.get("original_title"):
                changes.append(f"标题优化: '{opt_slide.get('original_title')}' -> '{opt_slide.get('new_title')}'")
            suggestions.extend(opt_slide.get("suggestions", []))

        # 保存（当前版本保存原文件副本）
        prs.save(str(output_path))

        return PolishResult(
            original_path=str(pptx_path),
            output_path=str(output_path),
            changes=changes,
            suggestions=suggestions,
        )

    def apply_color_scheme(
        self,
        pptx_path: Union[str, Path],
        scheme_name: str = "academic_blue",
        output_path: Optional[Union[str, Path]] = None,
    ) -> str:
        """
        应用配色方案

        Args:
            pptx_path: PPT 文件路径
            scheme_name: 配色方案名称
            output_path: 输出路径

        Returns:
            输出文件路径
        """
        from pptx import Presentation
        from pptx.dml.color import RgbColor
        from pptx.util import Pt

        pptx_path = Path(pptx_path)
        if not output_path:
            output_path = pptx_path.parent / f"{pptx_path.stem}_{scheme_name}.pptx"

        scheme = self.COLOR_SCHEMES.get(scheme_name, self.COLOR_SCHEMES["academic_blue"])

        prs = Presentation(str(pptx_path))

        for slide in prs.slides:
            for shape in slide.shapes:
                # 处理形状填充
                if shape.has_fill and hasattr(shape.fill, 'solid'):
                    try:
                        if shape.fill.type is not None:
                            # 根据位置判断是标题区域还是内容区域
                            if shape.top and shape.top.inches < 1.5:
                                shape.fill.solid()
                                shape.fill.fore_color.rgb = RgbColor(*scheme["primary"])
                    except:
                        pass

                # 处理文字颜色
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            # 白色背景上的文字用深色
                            if shape.top and shape.top.inches < 1.5:
                                run.font.color.rgb = RgbColor(255, 255, 255)
                            else:
                                run.font.color.rgb = RgbColor(*scheme["text"])

        prs.save(str(output_path))
        logger.info(f"已应用配色方案 '{scheme['name']}': {output_path}")
        return str(output_path)

    def unify_fonts(
        self,
        pptx_path: Union[str, Path],
        font_scheme: str = "professional",
        output_path: Optional[Union[str, Path]] = None,
    ) -> str:
        """
        统一字体

        Args:
            pptx_path: PPT 文件路径
            font_scheme: 字体方案名称
            output_path: 输出路径

        Returns:
            输出文件路径
        """
        from pptx import Presentation
        from pptx.util import Pt

        pptx_path = Path(pptx_path)
        if not output_path:
            output_path = pptx_path.parent / f"{pptx_path.stem}_unified.pptx"

        scheme = self.FONT_SCHEMES.get(font_scheme, self.FONT_SCHEMES["professional"])

        prs = Presentation(str(pptx_path))

        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue

                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        # 判断是标题还是正文
                        is_title = (
                            shape.top and shape.top.inches < 1.5 or
                            (run.font.size and run.font.size.pt > 24)
                        )

                        if is_title:
                            run.font.name = scheme["title_font"]
                            run.font.size = Pt(scheme["title_size"])
                        else:
                            run.font.name = scheme["body_font"]
                            run.font.size = Pt(scheme["body_size"])

        prs.save(str(output_path))
        logger.info(f"已统一字体 '{scheme['name']}': {output_path}")
        return str(output_path)

    def add_page_numbers(
        self,
        pptx_path: Union[str, Path],
        output_path: Optional[Union[str, Path]] = None,
        start_from: int = 1,
        skip_title: bool = True,
    ) -> str:
        """
        添加页码

        Args:
            pptx_path: PPT 文件路径
            output_path: 输出路径
            start_from: 起始页码
            skip_title: 是否跳过标题页

        Returns:
            输出文件路径
        """
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        from pptx.dml.color import RgbColor

        pptx_path = Path(pptx_path)
        if not output_path:
            output_path = pptx_path.parent / f"{pptx_path.stem}_numbered.pptx"

        prs = Presentation(str(pptx_path))
        total_slides = len(prs.slides)

        page_num = start_from
        for i, slide in enumerate(prs.slides):
            if skip_title and i == 0:
                continue

            # 添加页码文本框
            num_box = slide.shapes.add_textbox(
                Inches(12.5), Inches(7), Inches(0.8), Inches(0.4)
            )
            tf = num_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"{page_num}/{total_slides - (1 if skip_title else 0)}"
            p.font.size = Pt(12)
            p.font.color.rgb = RgbColor(128, 128, 128)
            p.alignment = PP_ALIGN.RIGHT

            page_num += 1

        prs.save(str(output_path))
        logger.info(f"已添加页码: {output_path}")
        return str(output_path)

    def polish(
        self,
        pptx_path: Union[str, Path],
        output_path: Optional[Union[str, Path]] = None,
        mode: PolishMode = PolishMode.FULL,
        color_scheme: str = "academic_blue",
        font_scheme: str = "professional",
        add_numbers: bool = True,
    ) -> PolishResult:
        """
        一键美化 PPT

        Args:
            pptx_path: 输入 PPT 路径
            output_path: 输出路径
            mode: 美化模式
            color_scheme: 配色方案
            font_scheme: 字体方案
            add_numbers: 是否添加页码

        Returns:
            PolishResult
        """
        pptx_path = Path(pptx_path)
        if not output_path:
            output_path = pptx_path.parent / f"{pptx_path.stem}_polished.pptx"

        changes = []
        suggestions = []

        # 复制原文件
        import shutil
        shutil.copy(str(pptx_path), str(output_path))
        current_path = output_path

        # 根据模式执行优化
        if mode in [PolishMode.CONTENT, PolishMode.FULL]:
            if self.client:
                result = self.polish_content(current_path, current_path)
                changes.extend(result.changes)
                suggestions.extend(result.suggestions)
            else:
                suggestions.append("未配置 LLM API，跳过内容优化")

        if mode in [PolishMode.STYLE, PolishMode.FULL]:
            self.apply_color_scheme(current_path, color_scheme, current_path)
            changes.append(f"应用配色方案: {color_scheme}")

            self.unify_fonts(current_path, font_scheme, current_path)
            changes.append(f"统一字体: {font_scheme}")

        if add_numbers:
            self.add_page_numbers(current_path, current_path)
            changes.append("添加页码")

        logger.info(f"PPT 美化完成: {output_path}")

        return PolishResult(
            original_path=str(pptx_path),
            output_path=str(output_path),
            changes=changes,
            suggestions=suggestions,
        )

    @staticmethod
    def list_color_schemes() -> Dict[str, str]:
        """列出可用配色方案"""
        return {k: v["name"] for k, v in PPTPolish.COLOR_SCHEMES.items()}

    @staticmethod
    def list_font_schemes() -> Dict[str, str]:
        """列出可用字体方案"""
        return {k: v["name"] for k, v in PPTPolish.FONT_SCHEMES.items()}
