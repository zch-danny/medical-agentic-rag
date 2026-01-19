"""
Paper2PPT - 从论文生成完整 PPT 演示文稿

分析论文结构，自动生成包含标题、摘要、方法、实验、结论的演示文稿
"""

import json
import os
import re
from dataclasses import dataclass, field
from enum import Enum
from pathlib import Path
from typing import Dict, List, Optional, Union

from loguru import logger
from openai import OpenAI
from tenacity import retry, stop_after_attempt, wait_exponential


class PPTStyle(str, Enum):
    """PPT 风格"""
    ACADEMIC = "academic"       # 学术风格（简洁专业）
    BUSINESS = "business"       # 商务风格
    MODERN = "modern"          # 现代简约
    COLORFUL = "colorful"      # 多彩活泼


@dataclass
class SlideContent:
    """单页幻灯片内容"""
    title: str
    content: List[str]          # 要点列表
    notes: Optional[str] = None # 演讲备注
    slide_type: str = "content" # title, content, section, conclusion
    image_prompt: Optional[str] = None  # 图片生成提示（可选）


@dataclass
class PPTContent:
    """PPT 完整内容"""
    title: str
    subtitle: Optional[str] = None
    author: Optional[str] = None
    slides: List[SlideContent] = field(default_factory=list)
    metadata: Dict = field(default_factory=dict)


class Paper2PPT:
    """
    论文转 PPT 生成器

    从论文内容自动生成结构化的演示文稿
    """

    ANALYSIS_PROMPT = """你是一个专业的学术演示文稿设计师。请分析以下论文内容，生成一份完整的 PPT 大纲。

要求：
1. 分析论文结构，识别关键章节（摘要、引言、方法、实验、结论等）
2. 为每个章节生成 1-3 页幻灯片
3. 每页幻灯片包含：标题、3-5 个要点、演讲备注
4. 要点应简洁有力，每条不超过 20 字
5. 总页数控制在 10-15 页

论文内容：
{content}

请以 JSON 格式输出，结构如下：
```json
{{
    "title": "演示文稿标题",
    "subtitle": "副标题（可选）",
    "author": "作者（如果论文中有）",
    "slides": [
        {{
            "slide_type": "title",
            "title": "标题页标题",
            "content": ["副标题或作者信息"],
            "notes": "开场白建议"
        }},
        {{
            "slide_type": "section",
            "title": "章节标题",
            "content": [],
            "notes": "章节过渡语"
        }},
        {{
            "slide_type": "content",
            "title": "页面标题",
            "content": ["要点1", "要点2", "要点3"],
            "notes": "详细讲解要点"
        }},
        {{
            "slide_type": "conclusion",
            "title": "总结",
            "content": ["结论1", "结论2"],
            "notes": "总结陈词"
        }}
    ]
}}
```

只输出 JSON，不要其他内容。"""

    ENHANCE_PROMPT = """请优化以下 PPT 内容，使其更加专业和吸引人：

原始内容：
{content}

优化要求：
1. 标题更加简洁有力
2. 要点更加精炼，使用动词开头
3. 添加过渡语和衔接
4. 补充演讲备注

输出格式与输入相同的 JSON 结构。只输出 JSON。"""

    def __init__(
        self,
        api_key: Optional[str] = None,
        base_url: Optional[str] = None,
        model: Optional[str] = None,
    ):
        self.api_key = api_key or os.getenv("LLM_API_KEY")
        self.base_url = base_url or os.getenv("LLM_BASE_URL", "https://api.deepseek.com")
        self.model = model or os.getenv("LLM_MODEL", "deepseek-chat")

        if not self.api_key:
            raise ValueError("LLM_API_KEY 未设置")

        self.client = OpenAI(api_key=self.api_key, base_url=self.base_url)
        logger.info(f"Paper2PPT 初始化完成: {self.base_url}, model={self.model}")

    def _extract_json(self, response: str) -> dict:
        """从 LLM 响应中提取 JSON"""
        # 尝试直接解析
        try:
            return json.loads(response)
        except json.JSONDecodeError:
            pass

        # 尝试提取 ```json ``` 块
        pattern = r'```json\s*([\s\S]*?)```'
        match = re.search(pattern, response)
        if match:
            try:
                return json.loads(match.group(1))
            except json.JSONDecodeError:
                pass

        # 尝试提取 { } 块
        pattern = r'\{[\s\S]*\}'
        match = re.search(pattern, response)
        if match:
            try:
                return json.loads(match.group(0))
            except json.JSONDecodeError:
                pass

        raise ValueError("无法从响应中提取有效的 JSON")

    @retry(
        stop=stop_after_attempt(3),
        wait=wait_exponential(multiplier=1, min=2, max=10),
        reraise=True,
    )
    def _call_llm(self, prompt: str) -> str:
        """调用 LLM"""
        response = self.client.chat.completions.create(
            model=self.model,
            messages=[
                {"role": "system", "content": "你是一个专业的学术演示文稿设计师，擅长将论文内容转化为清晰的 PPT。"},
                {"role": "user", "content": prompt},
            ],
            temperature=0.3,
            max_tokens=4096,
        )
        return response.choices[0].message.content

    def analyze(self, content: str, title: Optional[str] = None) -> PPTContent:
        """
        分析论文内容，生成 PPT 结构

        Args:
            content: 论文文本内容
            title: PPT 标题（可选，自动从论文提取）

        Returns:
            PPTContent: PPT 内容结构
        """
        # 截断过长内容
        max_length = 10000
        if len(content) > max_length:
            logger.warning(f"内容过长 ({len(content)} 字符)，截断到 {max_length}")
            content = content[:max_length] + "\n...[内容已截断]..."

        # 调用 LLM 分析
        prompt = self.ANALYSIS_PROMPT.format(content=content)
        logger.info("正在分析论文内容...")
        response = self._call_llm(prompt)

        # 解析结果
        data = self._extract_json(response)

        # 构建 PPTContent
        slides = []
        for slide_data in data.get("slides", []):
            slides.append(SlideContent(
                title=slide_data.get("title", ""),
                content=slide_data.get("content", []),
                notes=slide_data.get("notes"),
                slide_type=slide_data.get("slide_type", "content"),
                image_prompt=slide_data.get("image_prompt"),
            ))

        return PPTContent(
            title=title or data.get("title", "演示文稿"),
            subtitle=data.get("subtitle"),
            author=data.get("author"),
            slides=slides,
            metadata={"model": self.model, "content_length": len(content)},
        )

    def generate_pptx(
        self,
        ppt_content: PPTContent,
        output_path: Union[str, Path],
        style: PPTStyle = PPTStyle.ACADEMIC,
    ) -> str:
        """
        生成 PPTX 文件

        Args:
            ppt_content: PPT 内容
            output_path: 输出文件路径
            style: PPT 风格

        Returns:
            生成的文件路径
        """
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.dml.color import RgbColor
            from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
            from pptx.enum.shapes import MSO_SHAPE
        except ImportError:
            raise ImportError("python-pptx 未安装，请运行: pip install python-pptx")

        output_path = Path(output_path)

        # 创建演示文稿
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # 颜色方案
        color_schemes = {
            PPTStyle.ACADEMIC: {
                "primary": RgbColor(0x1a, 0x3a, 0x5c),     # 深蓝
                "secondary": RgbColor(0x4a, 0x90, 0xd9),   # 浅蓝
                "accent": RgbColor(0xe6, 0x7e, 0x22),      # 橙色
                "bg": RgbColor(0xff, 0xff, 0xff),          # 白色
                "text": RgbColor(0x2c, 0x3e, 0x50),        # 深灰
            },
            PPTStyle.MODERN: {
                "primary": RgbColor(0x2d, 0x3a, 0x4a),
                "secondary": RgbColor(0x00, 0xb8, 0x94),
                "accent": RgbColor(0xff, 0x6b, 0x6b),
                "bg": RgbColor(0xf8, 0xf9, 0xfa),
                "text": RgbColor(0x34, 0x3a, 0x40),
            },
            PPTStyle.COLORFUL: {
                "primary": RgbColor(0x66, 0x5c, 0xff),
                "secondary": RgbColor(0x00, 0xc9, 0xa7),
                "accent": RgbColor(0xff, 0x66, 0x90),
                "bg": RgbColor(0xff, 0xff, 0xff),
                "text": RgbColor(0x33, 0x33, 0x33),
            },
            PPTStyle.BUSINESS: {
                "primary": RgbColor(0x0d, 0x47, 0xa1),
                "secondary": RgbColor(0x42, 0xa5, 0xf5),
                "accent": RgbColor(0xff, 0xb3, 0x00),
                "bg": RgbColor(0xff, 0xff, 0xff),
                "text": RgbColor(0x21, 0x21, 0x21),
            },
        }
        colors = color_schemes.get(style, color_schemes[PPTStyle.ACADEMIC])

        for slide_content in ppt_content.slides:
            slide_layout = prs.slide_layouts[6]  # 空白布局
            slide = prs.slides.add_slide(slide_layout)

            if slide_content.slide_type == "title":
                # 标题页
                self._add_title_slide(slide, ppt_content, colors)

            elif slide_content.slide_type == "section":
                # 章节分隔页
                self._add_section_slide(slide, slide_content, colors)

            elif slide_content.slide_type == "conclusion":
                # 结论页
                self._add_conclusion_slide(slide, slide_content, colors)

            else:
                # 内容页
                self._add_content_slide(slide, slide_content, colors)

            # 添加演讲备注
            if slide_content.notes:
                notes_slide = slide.notes_slide
                notes_slide.notes_text_frame.text = slide_content.notes

        # 保存
        prs.save(str(output_path))
        logger.info(f"PPT 已生成: {output_path}")
        return str(output_path)

    def _add_title_slide(self, slide, ppt_content: PPTContent, colors: dict):
        """添加标题页"""
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN

        # 背景装饰条
        shape = slide.shapes.add_shape(
            1,  # MSO_SHAPE.RECTANGLE
            Inches(0), Inches(5.5), Inches(13.333), Inches(2)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = colors["primary"]
        shape.line.fill.background()

        # 主标题
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(12.333), Inches(1.5))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = ppt_content.title
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = colors["primary"]
        p.alignment = PP_ALIGN.CENTER

        # 副标题
        if ppt_content.subtitle:
            sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(12.333), Inches(0.8))
            tf = sub_box.text_frame
            p = tf.paragraphs[0]
            p.text = ppt_content.subtitle
            p.font.size = Pt(24)
            p.font.color.rgb = colors["text"]
            p.alignment = PP_ALIGN.CENTER

        # 作者
        if ppt_content.author:
            author_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.8), Inches(12.333), Inches(0.6))
            tf = author_box.text_frame
            p = tf.paragraphs[0]
            p.text = ppt_content.author
            p.font.size = Pt(18)
            p.font.color.rgb = RgbColor(0xff, 0xff, 0xff)
            p.alignment = PP_ALIGN.CENTER

    def _add_section_slide(self, slide, content: SlideContent, colors: dict):
        """添加章节分隔页"""
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN

        # 背景
        shape = slide.shapes.add_shape(
            1, Inches(0), Inches(0), Inches(13.333), Inches(7.5)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = colors["primary"]
        shape.line.fill.background()

        # 章节标题
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(12.333), Inches(1.5))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = content.title
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = RgbColor(0xff, 0xff, 0xff)
        p.alignment = PP_ALIGN.CENTER

    def _add_content_slide(self, slide, content: SlideContent, colors: dict):
        """添加内容页"""
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN

        # 顶部装饰条
        shape = slide.shapes.add_shape(
            1, Inches(0), Inches(0), Inches(13.333), Inches(1.2)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = colors["primary"]
        shape.line.fill.background()

        # 标题
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12.333), Inches(0.7))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = content.title
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = RgbColor(0xff, 0xff, 0xff)

        # 内容要点
        content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.733), Inches(5.5))
        tf = content_box.text_frame
        tf.word_wrap = True

        for i, point in enumerate(content.content):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            p.text = f"• {point}"
            p.font.size = Pt(22)
            p.font.color.rgb = colors["text"]
            p.space_after = Pt(16)
            p.level = 0

    def _add_conclusion_slide(self, slide, content: SlideContent, colors: dict):
        """添加结论页"""
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN

        # 背景渐变效果（用两个矩形模拟）
        shape1 = slide.shapes.add_shape(
            1, Inches(0), Inches(0), Inches(13.333), Inches(3.75)
        )
        shape1.fill.solid()
        shape1.fill.fore_color.rgb = colors["primary"]
        shape1.line.fill.background()

        shape2 = slide.shapes.add_shape(
            1, Inches(0), Inches(3.75), Inches(13.333), Inches(3.75)
        )
        shape2.fill.solid()
        shape2.fill.fore_color.rgb = colors["secondary"]
        shape2.line.fill.background()

        # 标题
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(1))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = content.title
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = RgbColor(0xff, 0xff, 0xff)
        p.alignment = PP_ALIGN.CENTER

        # 结论要点
        content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11.333), Inches(5))
        tf = content_box.text_frame
        tf.word_wrap = True

        for i, point in enumerate(content.content):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            p.text = f"✓ {point}"
            p.font.size = Pt(24)
            p.font.color.rgb = RgbColor(0xff, 0xff, 0xff)
            p.space_after = Pt(20)

    def generate(
        self,
        content: str,
        output_path: Union[str, Path],
        title: Optional[str] = None,
        style: PPTStyle = PPTStyle.ACADEMIC,
    ) -> str:
        """
        一键生成 PPT

        Args:
            content: 论文内容
            output_path: 输出路径
            title: PPT 标题
            style: PPT 风格

        Returns:
            生成的文件路径
        """
        # 分析内容
        ppt_content = self.analyze(content, title)

        # 生成 PPTX
        return self.generate_pptx(ppt_content, output_path, style)

    def generate_from_pdf(
        self,
        pdf_path: Union[str, Path],
        output_path: Optional[Union[str, Path]] = None,
        style: PPTStyle = PPTStyle.ACADEMIC,
    ) -> str:
        """
        从 PDF 生成 PPT

        Args:
            pdf_path: PDF 文件路径
            output_path: 输出路径（默认与 PDF 同目录）
            style: PPT 风格

        Returns:
            生成的文件路径
        """
        pdf_path = Path(pdf_path)
        if not pdf_path.exists():
            raise FileNotFoundError(f"PDF 文件不存在: {pdf_path}")

        # 确定输出路径
        if not output_path:
            output_path = pdf_path.with_suffix(".pptx")

        # 加载 PDF
        from src.document_loader import MinerUDocumentLoader
        loader = MinerUDocumentLoader()
        chunks = loader.load(str(pdf_path), chunk_size=10000, chunk_overlap=0)

        if not chunks:
            raise ValueError(f"无法从 PDF 提取文本: {pdf_path}")

        content = "\n\n".join(c["text"] for c in chunks)

        return self.generate(content, output_path, title=pdf_path.stem, style=style)
