"""
Paper2Figure 渲染模块

将 Mermaid 代码渲染为 SVG 和 PPTX 文件
"""

import base64
import subprocess
import tempfile
from io import BytesIO
from pathlib import Path
from typing import Optional, Union

from loguru import logger

from .core import FigureResult


class FigureRenderer:
    """
    图表渲染器

    支持将 Mermaid 代码渲染为：
    - SVG 图片（通过 mermaid-cli 或在线服务）
    - PPTX 演示文稿
    - HTML 预览页面
    """

    # Mermaid 在线渲染服务
    MERMAID_LIVE_URL = "https://mermaid.ink/svg"

    def __init__(self, output_dir: Optional[Union[str, Path]] = None):
        """
        初始化渲染器

        Args:
            output_dir: 输出目录（默认为临时目录）
        """
        if output_dir:
            self.output_dir = Path(output_dir)
            self.output_dir.mkdir(parents=True, exist_ok=True)
        else:
            self.output_dir = Path(tempfile.gettempdir()) / "paper2figure"
            self.output_dir.mkdir(parents=True, exist_ok=True)

        # 检查 mermaid-cli 是否可用
        self.mmdc_available = self._check_mmdc()

    def _check_mmdc(self) -> bool:
        """检查 mermaid-cli (mmdc) 是否可用"""
        try:
            result = subprocess.run(
                ["mmdc", "--version"],
                capture_output=True,
                text=True,
                timeout=5,
            )
            if result.returncode == 0:
                logger.debug(f"mermaid-cli 可用: {result.stdout.strip()}")
                return True
        except (subprocess.TimeoutExpired, FileNotFoundError):
            pass

        logger.debug("mermaid-cli 不可用，将使用在线渲染服务")
        return False

    def _render_with_mmdc(self, mermaid_code: str, output_path: Path) -> bool:
        """使用 mermaid-cli 渲染"""
        try:
            # 创建临时 mermaid 文件
            temp_mmd = self.output_dir / "temp.mmd"
            temp_mmd.write_text(mermaid_code, encoding="utf-8")

            result = subprocess.run(
                ["mmdc", "-i", str(temp_mmd), "-o", str(output_path), "-b", "transparent"],
                capture_output=True,
                text=True,
                timeout=30,
            )

            temp_mmd.unlink(missing_ok=True)

            if result.returncode == 0:
                return True
            else:
                logger.error(f"mmdc 渲染失败: {result.stderr}")
                return False

        except Exception as e:
            logger.error(f"mmdc 渲染异常: {e}")
            return False

    def _render_with_online(self, mermaid_code: str) -> Optional[str]:
        """使用在线服务渲染 Mermaid 为 SVG"""
        try:
            import requests

            # 使用 mermaid.ink 服务
            # URL 格式: https://mermaid.ink/svg/{base64编码的mermaid代码}
            encoded = base64.urlsafe_b64encode(mermaid_code.encode()).decode()
            url = f"{self.MERMAID_LIVE_URL}/{encoded}"

            response = requests.get(url, timeout=30)
            if response.status_code == 200:
                return response.text
            else:
                logger.error(f"在线渲染失败: HTTP {response.status_code}")
                return None

        except ImportError:
            logger.warning("requests 未安装，无法使用在线渲染")
            return None
        except Exception as e:
            logger.error(f"在线渲染异常: {e}")
            return None

    def render_svg(self, result: FigureResult, output_path: Optional[Union[str, Path]] = None) -> Optional[str]:
        """
        将 Mermaid 代码渲染为 SVG

        Args:
            result: FigureResult 对象
            output_path: 输出文件路径（可选）

        Returns:
            SVG 代码字符串，失败返回 None
        """
        mermaid_code = result.mermaid_code

        # 方式1：使用本地 mermaid-cli
        if self.mmdc_available and output_path:
            output_path = Path(output_path)
            if self._render_with_mmdc(mermaid_code, output_path):
                svg_content = output_path.read_text(encoding="utf-8")
                result.svg_code = svg_content
                return svg_content

        # 方式2：使用在线服务
        svg_content = self._render_with_online(mermaid_code)
        if svg_content:
            result.svg_code = svg_content

            # 保存到文件
            if output_path:
                output_path = Path(output_path)
                output_path.write_text(svg_content, encoding="utf-8")

            return svg_content

        logger.warning("SVG 渲染失败，返回原始 Mermaid 代码")
        return None

    def render_pptx(
        self,
        result: FigureResult,
        output_path: Optional[Union[str, Path]] = None,
        template_path: Optional[Union[str, Path]] = None,
    ) -> Optional[str]:
        """
        将图表导出为 PPTX 文件

        Args:
            result: FigureResult 对象
            output_path: 输出文件路径
            template_path: PPT 模板路径（可选）

        Returns:
            输出文件路径，失败返回 None
        """
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.dml.color import RgbColor
            from pptx.enum.text import PP_ALIGN
        except ImportError:
            logger.error("python-pptx 未安装，请运行: pip install python-pptx")
            return None

        # 确定输出路径
        if output_path:
            output_path = Path(output_path)
        else:
            safe_title = "".join(c for c in result.title if c.isalnum() or c in " _-")[:50]
            output_path = self.output_dir / f"{safe_title}.pptx"

        try:
            # 创建演示文稿
            if template_path and Path(template_path).exists():
                prs = Presentation(str(template_path))
            else:
                prs = Presentation()

            # 设置幻灯片大小为 16:9
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)

            # 添加标题幻灯片
            slide_layout = prs.slide_layouts[6]  # 空白布局
            slide = prs.slides.add_slide(slide_layout)

            # 添加标题
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(1))
            title_frame = title_box.text_frame
            title_para = title_frame.paragraphs[0]
            title_para.text = result.title
            title_para.font.size = Pt(32)
            title_para.font.bold = True
            title_para.alignment = PP_ALIGN.CENTER

            # 添加 Mermaid 代码区域（作为占位）
            # 由于 PPTX 无法直接渲染 Mermaid，我们添加代码和说明
            code_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.333), Inches(5))
            code_frame = code_box.text_frame
            code_frame.word_wrap = True

            # 添加说明
            intro_para = code_frame.paragraphs[0]
            intro_para.text = "📊 图表代码 (Mermaid)"
            intro_para.font.size = Pt(18)
            intro_para.font.bold = True

            # 添加代码
            code_para = code_frame.add_paragraph()
            code_para.text = result.mermaid_code
            code_para.font.size = Pt(12)
            code_para.font.name = "Consolas"

            # 添加使用说明
            note_para = code_frame.add_paragraph()
            note_para.text = "\n💡 提示: 复制以上代码到 https://mermaid.live 在线预览和编辑"
            note_para.font.size = Pt(11)
            note_para.font.italic = True

            # 如果有 SVG，尝试添加图片
            if result.svg_code:
                try:
                    # 将 SVG 转换为 PNG（需要 cairosvg）
                    import cairosvg
                    png_data = cairosvg.svg2png(bytestring=result.svg_code.encode())

                    # 添加新幻灯片放图片
                    img_slide = prs.slides.add_slide(slide_layout)

                    # 添加标题
                    img_title = img_slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
                    img_title.text_frame.paragraphs[0].text = result.title
                    img_title.text_frame.paragraphs[0].font.size = Pt(28)
                    img_title.text_frame.paragraphs[0].font.bold = True
                    img_title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

                    # 添加图片
                    img_stream = BytesIO(png_data)
                    img_slide.shapes.add_picture(img_stream, Inches(1), Inches(1.5), width=Inches(11))

                except ImportError:
                    logger.debug("cairosvg 未安装，跳过 SVG 图片嵌入")
                except Exception as e:
                    logger.debug(f"SVG 转 PNG 失败: {e}")

            # 保存
            prs.save(str(output_path))
            result.pptx_path = str(output_path)
            logger.info(f"PPTX 已生成: {output_path}")
            return str(output_path)

        except Exception as e:
            logger.error(f"PPTX 生成失败: {e}")
            return None

    def render_html(self, result: FigureResult, output_path: Optional[Union[str, Path]] = None) -> str:
        """
        生成包含 Mermaid 图表的 HTML 页面

        Args:
            result: FigureResult 对象
            output_path: 输出文件路径（可选）

        Returns:
            HTML 内容字符串
        """
        html_template = """<!DOCTYPE html>
<html lang="zh-CN">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{title}</title>
    <script src="https://cdn.jsdelivr.net/npm/mermaid/dist/mermaid.min.js"></script>
    <style>
        body {{
            font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', Roboto, sans-serif;
            max-width: 1200px;
            margin: 0 auto;
            padding: 20px;
            background: #f5f5f5;
        }}
        h1 {{
            text-align: center;
            color: #333;
        }}
        .mermaid {{
            background: white;
            padding: 20px;
            border-radius: 8px;
            box-shadow: 0 2px 10px rgba(0,0,0,0.1);
        }}
        .code-block {{
            background: #2d2d2d;
            color: #f8f8f2;
            padding: 15px;
            border-radius: 8px;
            margin-top: 20px;
            overflow-x: auto;
        }}
        .code-block pre {{
            margin: 0;
            font-family: 'Consolas', 'Monaco', monospace;
            font-size: 14px;
        }}
        .description {{
            color: #666;
            text-align: center;
            margin-bottom: 20px;
        }}
        .actions {{
            text-align: center;
            margin-top: 20px;
        }}
        .actions a {{
            display: inline-block;
            padding: 10px 20px;
            background: #4CAF50;
            color: white;
            text-decoration: none;
            border-radius: 5px;
            margin: 5px;
        }}
        .actions a:hover {{
            background: #45a049;
        }}
    </style>
</head>
<body>
    <h1>{title}</h1>
    <p class="description">{description}</p>

    <div class="mermaid">
{mermaid_code}
    </div>

    <div class="actions">
        <a href="https://mermaid.live/edit#pako:{encoded}" target="_blank">📝 在线编辑</a>
        <a href="javascript:void(0)" onclick="copyCode()">📋 复制代码</a>
    </div>

    <h3>Mermaid 源代码</h3>
    <div class="code-block">
        <pre id="code">{mermaid_code_escaped}</pre>
    </div>

    <script>
        mermaid.initialize({{ startOnLoad: true, theme: 'default' }});

        function copyCode() {{
            const code = document.getElementById('code').textContent;
            navigator.clipboard.writeText(code).then(() => {{
                alert('代码已复制到剪贴板！');
            }});
        }}
    </script>
</body>
</html>"""

        import html
        import zlib

        # 编码用于 mermaid.live
        try:
            compressed = zlib.compress(result.mermaid_code.encode(), 9)
            encoded = base64.urlsafe_b64encode(compressed).decode()
        except:
            encoded = ""

        html_content = html_template.format(
            title=result.title,
            description=result.description,
            mermaid_code=result.mermaid_code,
            mermaid_code_escaped=html.escape(result.mermaid_code),
            encoded=encoded,
        )

        if output_path:
            output_path = Path(output_path)
            output_path.write_text(html_content, encoding="utf-8")
            logger.info(f"HTML 已生成: {output_path}")

        return html_content

    def render_all(
        self,
        result: FigureResult,
        output_dir: Optional[Union[str, Path]] = None,
        formats: Optional[list] = None,
    ) -> dict:
        """
        渲染所有格式

        Args:
            result: FigureResult 对象
            output_dir: 输出目录
            formats: 要生成的格式列表，默认 ["html", "pptx"]

        Returns:
            dict: {"html": path, "pptx": path, "svg": path}
        """
        if output_dir:
            output_dir = Path(output_dir)
        else:
            output_dir = self.output_dir

        output_dir.mkdir(parents=True, exist_ok=True)

        if formats is None:
            formats = ["html", "pptx"]

        safe_title = "".join(c for c in result.title if c.isalnum() or c in " _-")[:50] or "figure"
        outputs = {}

        if "svg" in formats:
            svg_path = output_dir / f"{safe_title}.svg"
            if self.render_svg(result, svg_path):
                outputs["svg"] = str(svg_path)

        if "html" in formats:
            html_path = output_dir / f"{safe_title}.html"
            self.render_html(result, html_path)
            outputs["html"] = str(html_path)

        if "pptx" in formats:
            pptx_path = output_dir / f"{safe_title}.pptx"
            if self.render_pptx(result, pptx_path):
                outputs["pptx"] = str(pptx_path)

        return outputs
